# src/email/attachment_handler.py
import asyncio
import logging
import mimetypes
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, BinaryIO
import base64
import hashlib
from docx import Document


# Office document processing
import docx
from openpyxl import load_workbook
from pptx import Presentation
import PyPDF2
import pandas as pd

# Image processing
from PIL import Image
import io

from src.ai.model_manager import ai_manager
from src.utils.cache_manager import cache_manager
from src.utils.text_utils import clean_text, extract_keywords

logger = logging.getLogger(__name__)

class AttachmentType:
    """Supported attachment types"""
    PDF = "pdf"
    WORD = "word"
    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    IMAGE = "image"
    TEXT = "text"
    ARCHIVE = "archive"
    UNKNOWN = "unknown"

class AttachmentInfo:
    """Attachment metadata and content"""
    def __init__(self, filename: str, file_type: str, size: int, mime_type: str):
        self.filename = filename
        self.file_type = file_type
        self.size = size
        self.mime_type = mime_type
        self.content_text = ""
        self.metadata = {}
        self.summary = ""
        self.key_points = []
        self.requires_action = False
        self.confidence = 0.0
        self.processing_error = None

class AttachmentProcessor:
    """Advanced attachment processing with AI analysis"""
    
    def __init__(self):
        self.max_file_size = 50 * 1024 * 1024  # 50MB limit
        self.supported_types = {
            '.pdf': AttachmentType.PDF,
            '.doc': AttachmentType.WORD,
            '.docx': AttachmentType.WORD,
            '.xls': AttachmentType.EXCEL,
            '.xlsx': AttachmentType.EXCEL,
            '.xlsm': AttachmentType.EXCEL,
            '.ppt': AttachmentType.POWERPOINT,
            '.pptx': AttachmentType.POWERPOINT,
            '.txt': AttachmentType.TEXT,
            '.csv': AttachmentType.TEXT,
            '.jpg': AttachmentType.IMAGE,
            '.jpeg': AttachmentType.IMAGE,
            '.png': AttachmentType.IMAGE,
            '.gif': AttachmentType.IMAGE,
            '.zip': AttachmentType.ARCHIVE,
            '.rar': AttachmentType.ARCHIVE,
        }
    
    async def process_attachment(self, attachment_data: bytes, filename: str, 
                               mime_type: str = None) -> AttachmentInfo:
        """Process attachment and extract meaningful information"""
        try:
            # Validate file size
            if len(attachment_data) > self.max_file_size:
                raise ValueError(f"File too large: {len(attachment_data)} bytes")
            
            # Determine file type
            file_extension = Path(filename).suffix.lower()
            file_type = self.supported_types.get(file_extension, AttachmentType.UNKNOWN)
            
            if not mime_type:
                mime_type, _ = mimetypes.guess_type(filename)
            
            # Create attachment info
            attachment_info = AttachmentInfo(filename, file_type, len(attachment_data), mime_type)
            
            # Process based on file type
            if file_type == AttachmentType.PDF:
                await self._process_pdf(attachment_data, attachment_info)
            elif file_type == AttachmentType.WORD:
                await self._process_word_document(attachment_data, attachment_info)
            elif file_type == AttachmentType.EXCEL:
                await self._process_excel_document(attachment_data, attachment_info)
            elif file_type == AttachmentType.POWERPOINT:
                await self._process_powerpoint(attachment_data, attachment_info)
            elif file_type == AttachmentType.TEXT:
                await self._process_text_file(attachment_data, attachment_info)
            elif file_type == AttachmentType.IMAGE:
                await self._process_image(attachment_data, attachment_info)
            elif file_type == AttachmentType.ARCHIVE:
                await self._process_archive(attachment_data, attachment_info)
            else:
                attachment_info.processing_error = f"Unsupported file type: {file_extension}"
                attachment_info.confidence = 0.1
            
            # Generate AI analysis if content was extracted
            if attachment_info.content_text and not attachment_info.processing_error:
                await self._analyze_content_with_ai(attachment_info)
            
            return attachment_info
            
        except Exception as e:
            logger.error(f"Attachment processing failed for {filename}: {str(e)}")
            error_info = AttachmentInfo(filename, AttachmentType.UNKNOWN, len(attachment_data), mime_type)
            error_info.processing_error = str(e)
            error_info.confidence = 0.0
            return error_info
    
    async def _process_pdf(self, pdf_data: bytes, info: AttachmentInfo):
        """Extract text and metadata from PDF"""
        try:
            pdf_stream = io.BytesIO(pdf_data)
            pdf_reader = PyPDF2.PdfReader(pdf_stream)
            
            # Extract metadata
            info.metadata = {
                'pages': len(pdf_reader.pages),
                'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                'creation_date': str(pdf_reader.metadata.get('/CreationDate', '')) if pdf_reader.metadata else '',
            }
            
            # Extract text from all pages
            text_content = []
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"Page {page_num + 1}:\n{page_text}")
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num + 1}: {str(e)}")
            
            info.content_text = "\n\n".join(text_content)
            info.confidence = 0.8 if info.content_text else 0.3
            
        except Exception as e:
            info.processing_error = f"PDF processing error: {str(e)}"
            info.confidence = 0.1
    
    async def _process_word_document(self, doc_data: bytes, info: AttachmentInfo):
        """Extract text and metadata from Word document"""
        try:
            doc_stream = io.BytesIO(doc_data)
            doc = docx.Document(doc_stream)
            
            # Extract text from paragraphs
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            # Extract text from tables
            table_content = []
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                table_content.append("\n".join(table_text))
            
            # Combine all content
            all_content = paragraphs
            if table_content:
                all_content.append("\n--- Tables ---\n")
                all_content.extend(table_content)
            
            info.content_text = "\n\n".join(all_content)
            
            # Extract metadata
            core_props = doc.core_properties
            info.metadata = {
                'author': core_props.author or '',
                'title': core_props.title or '',
                'subject': core_props.subject or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'word_count': len(info.content_text.split()),
                'paragraphs': len(paragraphs),
                'tables': len(doc.tables)
            }
            
            info.confidence = 0.9 if info.content_text else 0.3
            
        except Exception as e:
            info.processing_error = f"Word document processing error: {str(e)}"
            info.confidence = 0.1
    
    async def _process_excel_document(self, excel_data: bytes, info: AttachmentInfo):
        """Extract text and data from Excel document"""
        try:
            excel_stream = io.BytesIO(excel_data)
            workbook = load_workbook(excel_stream, data_only=True)
            
            sheet_contents = []
            total_rows = 0
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Convert sheet to DataFrame for easier processing
                data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):  # Skip empty rows
                        data.append([str(cell) if cell is not None else '' for cell in row])
                
                if data:
                    # Create summary of sheet content
                    df = pd.DataFrame(data[1:], columns=data[0] if data else None)
                    
                    sheet_summary = f"Sheet: {sheet_name}\n"
                    sheet_summary += f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
                    
                    if not df.empty:
                        # Add column names
                        sheet_summary += f"Columns: {', '.join(str(col) for col in df.columns)}\n"
                        
                        # Add sample data (first few rows)
                        sample_rows = min(5, len(df))
                        sheet_summary += f"Sample data (first {sample_rows} rows):\n"
                        sheet_summary += df.head(sample_rows).to_string(index=False)
                        
                        # Add data summary statistics for numeric columns
                        numeric_cols = df.select_dtypes(include=['number']).columns
                        if len(numeric_cols) > 0:
                            sheet_summary += f"\n\nNumeric summary:\n"
                            sheet_summary += df[numeric_cols].describe().to_string()
                    
                    sheet_contents.append(sheet_summary)
                    total_rows += len(df)
            
            info.content_text = "\n\n" + "="*50 + "\n\n".join(sheet_contents)
            
            # Extract metadata
            info.metadata = {
                'sheets': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames,
                'total_rows': total_rows,
                'workbook_properties': {
                    'created': str(workbook.properties.created) if workbook.properties.created else '',
                    'modified': str(workbook.properties.modified) if workbook.properties.modified else '',
                    'creator': workbook.properties.creator or '',
                }
            }
            
            info.confidence = 0.8 if sheet_contents else 0.3
            
        except Exception as e:
            info.processing_error = f"Excel processing error: {str(e)}"
            info.confidence = 0.1
    
    async def _process_powerpoint(self, ppt_data: bytes, info: AttachmentInfo):
        """Extract text from PowerPoint presentation"""
        try:
            ppt_stream = io.BytesIO(ppt_data)
            presentation = Presentation(ppt_stream)
            
            slide_contents = []
            total_text_boxes = 0
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                        total_text_boxes += 1
                
                if slide_text:
                    slide_content = f"Slide {slide_num}:\n" + "\n".join(slide_text)
                    slide_contents.append(slide_content)
            
            info.content_text = "\n\n" + "-"*40 + "\n\n".join(slide_contents)
            
            # Extract metadata
            info.metadata = {
                'slides': len(presentation.slides),
                'text_boxes': total_text_boxes,
                'slide_layout_info': [
                    f"Slide {i+1}: {len([s for s in slide.shapes if hasattr(s, 'text')])} text elements"
                    for i, slide in enumerate(presentation.slides)
                ]
            }
            
            info.confidence = 0.8 if slide_contents else 0.3
            
        except Exception as e:
            info.processing_error = f"PowerPoint processing error: {str(e)}"
            info.confidence = 0.1
    
    async def _process_text_file(self, text_data: bytes, info: AttachmentInfo):
        """Process plain text or CSV files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text_content = None
            
            for encoding in encodings:
                try:
                    text_content = text_data.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text_content:
                raise ValueError("Could not decode text file with any encoding")
            
            info.content_text = text_content
            
            # Special handling for CSV files
            if info.filename.lower().endswith('.csv'):
                try:
                    # Parse as CSV for better understanding
                    csv_stream = io.StringIO(text_content)
                    df = pd.read_csv(csv_stream)
                    
                    csv_summary = f"CSV Analysis:\n"
                    csv_summary += f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
                    csv_summary += f"Columns: {', '.join(df.columns)}\n"
                    
                    # Add sample data
                    sample_rows = min(5, len(df))
                    csv_summary += f"\nSample data (first {sample_rows} rows):\n"
                    csv_summary += df.head(sample_rows).to_string(index=False)
                    
                    # Add summary for numeric columns
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    if len(numeric_cols) > 0:
                        csv_summary += f"\n\nNumeric summary:\n"
                        csv_summary += df[numeric_cols].describe().to_string()
                    
                    info.content_text = csv_summary + "\n\n" + "="*50 + "\n\nRaw Content:\n" + text_content
                    
                    info.metadata = {
                        'file_type': 'csv',
                        'rows': len(df),
                        'columns': len(df.columns),
                        'column_names': list(df.columns)
                    }
                    
                except Exception:
                    # If CSV parsing fails, treat as regular text
                    pass
            
            # General text metadata
            if 'file_type' not in info.metadata:
                info.metadata = {
                    'file_type': 'text',
                    'lines': len(text_content.splitlines()),
                    'characters': len(text_content),
                    'words': len(text_content.split())
                }
            
            info.confidence = 0.9
            
        except Exception as e:
            info.processing_error = f"Text processing error: {str(e)}"
            info.confidence = 0.1
    
    async def _process_image(self, image_data: bytes, info: AttachmentInfo):
        """Extract metadata from images and optionally OCR text"""
        try:
            image_stream = io.BytesIO(image_data)
            image = Image.open(image_stream)
            
            # Extract basic image metadata
            info.metadata = {
                'format': image.format,
                'mode': image.mode,
                'size': image.size,
                'width': image.size[0],
                'height': image.size[1],
            }
            
            # Extract EXIF data if available
            if hasattr(image, '_getexif'):
                exif_data = image._getexif()
                if exif_data:
                    info.metadata['exif'] = {k: str(v) for k, v in exif_data.items() if k < 1000}  # Limit EXIF data
            
            # For now, basic image description
            info.content_text = f"Image file: {info.filename}\n"
            info.content_text += f"Format: {image.format}\n"
            info.content_text += f"Dimensions: {image.size[0]} x {image.size[1]}\n"
            info.content_text += f"Mode: {image.mode}"
            
            # TODO: Add OCR capability here if needed
            # This would require additional dependencies like pytesseract
            
            info.confidence = 0.6
            
        except Exception as e:
            info.processing_error = f"Image processing error: {str(e)}"
            info.confidence = 0.1
    
    async def _process_archive(self, archive_data: bytes, info: AttachmentInfo):
        """Extract file list from archive files"""
        try:
            archive_stream = io.BytesIO(archive_data)
            
            if info.filename.lower().endswith('.zip'):
                with zipfile.ZipFile(archive_stream, 'r') as zip_file:
                    file_list = zip_file.namelist()
                    
                    # Get file info
                    file_details = []
                    total_size = 0
                    
                    for filename in file_list:
                        file_info = zip_file.getinfo(filename)
                        file_details.append({
                            'name': filename,
                            'size': file_info.file_size,
                            'compressed_size': file_info.compress_size,
                            'date': f"{file_info.date_time}"
                        })
                        total_size += file_info.file_size
                    
                    # Create summary
                    info.content_text = f"Archive Contents ({len(file_list)} files):\n\n"
                    
                    # List all files
                    for detail in file_details:
                        size_kb = detail['size'] / 1024
                        info.content_text += f"• {detail['name']} ({size_kb:.1f} KB)\n"
                    
                    info.content_text += f"\nTotal uncompressed size: {total_size / 1024:.1f} KB"
                    
                    info.metadata = {
                        'archive_type': 'zip',
                        'file_count': len(file_list),
                        'total_uncompressed_size': total_size,
                        'files': file_details
                    }
            
            info.confidence = 0.7
            
        except Exception as e:
            info.processing_error = f"Archive processing error: {str(e)}"
            info.confidence = 0.1
    
    async def _analyze_content_with_ai(self, info: AttachmentInfo):
        """Analyze extracted content using AI for insights"""
        try:
            if not info.content_text or len(info.content_text) < 100:
                return
            
            # Truncate content for AI analysis (to manage token limits)
            analysis_content = info.content_text[:3000]
            
            analysis_prompt = f"""
            Analyze this document content and provide insights:
            
            Document: {info.filename}
            Type: {info.file_type}
            Content: {analysis_content}
            
            Please provide:
            1. Brief summary (2-3 sentences)
            2. Key points or highlights (up to 5 bullet points)
            3. Does this document require any action from the recipient? (yes/no and why)
            4. What type of document is this? (report, proposal, invoice, etc.)
            5. Any important dates, deadlines, or numbers mentioned?
            
            Respond in JSON format.
            """
            
            ai_response = await ai_manager.generate_response(analysis_prompt, temperature=0.3)
            
            # Try to parse JSON response
            import json
            import re
            json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
            
            if json_match:
                analysis = json.loads(json_match.group())
                
                info.summary = analysis.get('summary', 'AI analysis not available')
                info.key_points = analysis.get('key_points', [])
                
                # Check if action required
                action_analysis = analysis.get('action_required', '').lower()
                info.requires_action = 'yes' in action_analysis
                
                # Add additional metadata from AI
                if 'document_type' in analysis:
                    info.metadata['document_type'] = analysis['document_type']
                
                if 'important_dates' in analysis:
                    info.metadata['important_dates'] = analysis['important_dates']
                
                # Boost confidence if AI analysis was successful
                info.confidence = min(info.confidence + 0.2, 1.0)
            
        except Exception as e:
            logger.warning(f"AI content analysis failed for {info.filename}: {str(e)}")
            # Don't fail the entire processing if AI analysis fails
            info.summary = "AI analysis not available"
    
    async def batch_process_attachments(self, attachments: List[Tuple[bytes, str, str]]) -> List[AttachmentInfo]:
        """Process multiple attachments in parallel"""
        try:
            # Create processing tasks
            tasks = []
            for attachment_data, filename, mime_type in attachments:
                task = self.process_attachment(attachment_data, filename, mime_type)
                tasks.append(task)
            
            # Process in parallel with limited concurrency
            semaphore = asyncio.Semaphore(3)  # Limit to 3 concurrent processes
            
            async def process_with_semaphore(task):
                async with semaphore:
                    return await task
            
            results = await asyncio.gather(
                *[process_with_semaphore(task) for task in tasks],
                return_exceptions=True
            )
            
            # Filter out exceptions
            processed_attachments = []
            for i, result in enumerate(results):
                if isinstance(result, AttachmentInfo):
                    processed_attachments.append(result)
                else:
                    # Create error attachment info
                    filename = attachments[i][1] if i < len(attachments) else "unknown"
                    error_info = AttachmentInfo(filename, AttachmentType.UNKNOWN, 0, "unknown")
                    error_info.processing_error = str(result)
                    processed_attachments.append(error_info)
            
            return processed_attachments
            
        except Exception as e:
            logger.error(f"Batch attachment processing failed: {str(e)}")
            return []
    
    def get_attachment_summary(self, attachments: List[AttachmentInfo]) -> Dict[str, Any]:
        """Generate summary statistics for processed attachments"""
        if not attachments:
            return {}
        
        summary = {
            'total_attachments': len(attachments),
            'successfully_processed': len([a for a in attachments if not a.processing_error]),
            'failed_processing': len([a for a in attachments if a.processing_error]),
            'types': {},
            'total_size': sum(a.size for a in attachments),
            'requires_action_count': len([a for a in attachments if a.requires_action]),
            'average_confidence': sum(a.confidence for a in attachments) / len(attachments),
        }
        
        # Count by type
        for attachment in attachments:
            file_type = attachment.file_type
            summary['types'][file_type] = summary['types'].get(file_type, 0) + 1
        
        return summary
    
    async def search_attachments_content(self, attachments: List[AttachmentInfo], 
                                       search_query: str) -> List[AttachmentInfo]:
        """Search through attachment content"""
        matching_attachments = []
        search_terms = search_query.lower().split()
        
        for attachment in attachments:
            if attachment.content_text:
                content_lower = attachment.content_text.lower()
                
                # Check if any search terms match
                if any(term in content_lower for term in search_terms):
                    matching_attachments.append(attachment)
                
                # Also check filename and summary
                elif any(term in attachment.filename.lower() for term in search_terms):
                    matching_attachments.append(attachment)
                
                elif any(term in attachment.summary.lower() for term in search_terms):
                    matching_attachments.append(attachment)
        
        return matching_attachments

# Global attachment processor instance
attachment_processor = AttachmentProcessor()